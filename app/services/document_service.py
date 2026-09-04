"""
Document Service – Handles document ingestion and text extraction.

Supports: PDF (with OCR fallback), DOCX, PPTX, TXT, CSV
Uses Tesseract OCR for scanned/image-based PDF pages.
"""

import csv
import io
import logging
import os
import tempfile
import uuid
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

from app.config import get_settings

logger = logging.getLogger(__name__)

# Supported file extensions
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".csv"}


@dataclass
class DocumentResult:
    """Result of document text extraction."""
    doc_id: str
    filename: str
    text: str
    page_count: int
    file_type: str
    file_size_bytes: int
    extracted_at: datetime = field(default_factory=datetime.now)
    metadata: dict = field(default_factory=dict)


class DocumentService:
    """Service for extracting text from various document formats."""

    def __init__(self):
        settings = get_settings()
        self.upload_dir = Path(settings.UPLOAD_DIR)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.tesseract_cmd = settings.TESSERACT_CMD
        self.max_file_size = settings.MAX_FILE_SIZE_MB * 1024 * 1024  # bytes

    async def process_file(self, file) -> DocumentResult:
        """
        Process an uploaded file and extract text.

        Args:
            file: FastAPI UploadFile object

        Returns:
            DocumentResult with extracted text and metadata.

        Raises:
            ValueError: If file type is unsupported or file is too large.
        """
        filename = file.filename
        ext = Path(filename).suffix.lower()

        if ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: '{ext}'. "
                f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
            )

        # Read file content
        content = await file.read()
        file_size = len(content)

        if file_size > self.max_file_size:
            settings = get_settings()
            raise ValueError(
                f"File too large ({file_size / (1024*1024):.1f} MB). "
                f"Max allowed: {settings.MAX_FILE_SIZE_MB} MB"
            )

        # Save to uploads directory
        doc_id = str(uuid.uuid4())
        save_path = self.upload_dir / f"{doc_id}_{filename}"
        save_path.write_bytes(content)

        logger.info(f"Processing document: {filename} ({file_size} bytes, type={ext})")

        # Extract text based on file type
        try:
            if ext == ".pdf":
                text, page_count = self._extract_pdf(content)
            elif ext in (".docx", ".doc"):
                text, page_count = self._extract_docx(content)
            elif ext in (".pptx", ".ppt"):
                text, page_count = self._extract_pptx(content)
            elif ext == ".txt":
                text, page_count = self._extract_txt(content)
            elif ext == ".csv":
                text, page_count = self._extract_csv(content)
            else:
                raise ValueError(f"No handler for extension: {ext}")
        except Exception as e:
            # Clean up saved file on failure
            save_path.unlink(missing_ok=True)
            logger.error(f"Text extraction failed for {filename}: {e}")
            raise RuntimeError(f"Failed to extract text from {filename}: {str(e)}")

        if not text.strip():
            logger.warning(f"No text extracted from {filename}")

        logger.info(
            f"Extracted {len(text)} chars, {page_count} pages from {filename}"
        )

        return DocumentResult(
            doc_id=doc_id,
            filename=filename,
            text=text,
            page_count=page_count,
            file_type=ext,
            file_size_bytes=file_size,
            metadata={
                "saved_path": str(save_path),
                "original_filename": filename,
            }
        )

    # ──────────────────────────────────────────────
    # Format-Specific Extractors
    # ──────────────────────────────────────────────

    def _extract_pdf(self, content: bytes) -> tuple[str, int]:
        """Extract text from PDF, with OCR fallback for scanned pages."""
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(content))
        pages_text = []
        ocr_used = False

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""

            # If very little text extracted, try OCR
            if len(text.strip()) < 50:
                ocr_text = self._ocr_pdf_page(content, i)
                if ocr_text and len(ocr_text.strip()) > len(text.strip()):
                    text = ocr_text
                    ocr_used = True

            pages_text.append(f"[Page {i+1}]\n{text}")

        if ocr_used:
            logger.info("OCR was used for some pages")

        return "\n\n".join(pages_text), len(reader.pages)

    def _ocr_pdf_page(self, pdf_bytes: bytes, page_number: int) -> str:
        """OCR a single PDF page using PaddleOCR, then Tesseract."""
        try:
            from pdf2image import convert_from_bytes

            images = convert_from_bytes(
                pdf_bytes,
                first_page=page_number + 1,
                last_page=page_number + 1,
                dpi=300
            )
            if not images:
                return ""

            image = images[0]

            try:
                from paddleocr import PaddleOCR
                import numpy as np

                paddle_ocr = PaddleOCR(use_angle_cls=True, lang="en", show_log=False)
                result = paddle_ocr.ocr(np.array(image.convert("RGB")), cls=True)
                if result and result[0]:
                    lines = []
                    for line in result[0]:
                        if line and len(line) >= 2 and line[1]:
                            word_info = line[1]
                            if isinstance(word_info, (list, tuple)) and word_info:
                                lines.append(str(word_info[0]))
                    text = "\n".join(lines).strip()
                    if text:
                        return text
            except ImportError:
                logger.info("PaddleOCR not installed, trying Tesseract for PDF OCR")
            except Exception as e:
                logger.warning(f"PaddleOCR PDF page failed: {e}")

            try:
                import pytesseract

                pytesseract.pytesseract.tesseract_cmd = self.tesseract_cmd
                return pytesseract.image_to_string(image)
            except ImportError:
                logger.warning("pytesseract not available for OCR")
            except Exception as e:
                logger.warning(f"OCR failed for page {page_number}: {e}")

        except ImportError:
            logger.warning("pdf2image not available for OCR")
        except Exception as e:
            logger.warning(f"OCR failed for page {page_number}: {e}")

        return ""

    def _extract_docx(self, content: bytes) -> tuple[str, int]:
        """Extract text from DOCX files including tables."""
        from docx import Document

        doc = Document(io.BytesIO(content))
        parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)

        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(row_text))
            if table_text:
                parts.append("\n".join(table_text))

        # Approximate page count (roughly 3000 chars per page)
        full_text = "\n\n".join(parts)
        page_count = max(1, len(full_text) // 3000)

        return full_text, page_count

    def _extract_pptx(self, content: bytes) -> tuple[str, int]:
        """Extract text from PowerPoint presentations."""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides_text = []

        for i, slide in enumerate(prs.slides):
            slide_parts = [f"[Slide {i+1}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text)
            slides_text.append("\n".join(slide_parts))

        return "\n\n".join(slides_text), len(prs.slides)

    def _extract_txt(self, content: bytes) -> tuple[str, int]:
        """Extract text from plain text files."""
        # Try UTF-8 first, then fall back to latin-1
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        page_count = max(1, len(text) // 3000)
        return text, page_count

    def _extract_csv(self, content: bytes) -> tuple[str, int]:
        """Extract text from CSV files as structured text."""
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            rows.append(" | ".join(row))

        full_text = "\n".join(rows)
        page_count = max(1, len(rows) // 50)

        return full_text, page_count

    def delete_file(self, doc_id: str) -> bool:
        """Delete the uploaded file from disk."""
        for f in self.upload_dir.iterdir():
            if f.name.startswith(doc_id):
                f.unlink()
                logger.info(f"Deleted file: {f.name}")
                return True
        return False


# Singleton instance
document_service = DocumentService()
